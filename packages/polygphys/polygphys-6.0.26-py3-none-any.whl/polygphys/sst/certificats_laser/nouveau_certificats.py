# -*- coding: utf-8 -*-
"""Créer de nouveaux certificats laser au besoin."""

# Bibliothèques standard
import subprocess
import time

from pathlib import Path
from datetime import datetime as dt
from subprocess import run

# Bibliothèque PIPy
import schedule
import pptx
import keyring
import getpass

import sqlalchemy as sqla
import pandas as pd

# Bibliothèques maison
from polygphys.outils.base_de_donnees import BaseDeDonnées, BaseTableau
from polygphys.outils.base_de_donnees.dtypes import column
from polygphys.outils.reseau.msforms import MSFormConfig, MSForm
from polygphys.outils.reseau import DisqueRéseau, OneDrive
from polygphys.outils.reseau.courriel import Courriel

# Définitions de classes


class SSTLaserCertificatsConfig(MSFormConfig):

    def default(self):
        return (Path(__file__).parent / 'nouveau_certificat.cfg').open().read()


class Certificat:

    def __init__(self, modèle):
        self.modèle = modèle

    def màj(self, nom, matricule):
        self.cert = pptx.Presentation(self.modèle)
        for forme in self.cert.slides[0].shapes:
            if forme.has_text_frame:
                for par in forme.text_frame.paragraphs:
                    for ligne in par.runs:
                        if ligne.text == 'nom':
                            ligne.text = str(nom)
                        elif ligne.text == 'matricule':
                            ligne.text = str(matricule)
                        elif ligne.text.startswith('Date'):
                            date = dt.today()
                            ligne.text = f'Date: {date.year}-{date.month:02}'

    def enregistrer(self, fichier):
        self.cert.save(fichier)

        fichier_pdf = fichier.parent.parent / 'pdf' / fichier.name
        run(['unoconv',
             '-f',
             'pdf',
             '-o',
             str(fichier_pdf),
             str(fichier)])


class Participants(BaseTableau):
    colonnes_standard = (column('index', int, primary_key=True),
                         column('matricule', str),
                         column('nom', str),
                         column('courriel', str))

    def __init__(self, adresse: str, reflect: bool = True):
        table = 'personnes'

        metadata = sqla.MetaData()
        db = BaseDeDonnées(adresse, metadata)

        if reflect:
            moteur = db.create_engine()
            metadata.reflect(moteur)
        else:
            sqla.Table(table,
                       metadata,
                       *self.colonnes_standard)

        super().__init__(db, table)

    def ajouter(self, matricule, nom, courriel) -> int:
        df = pd.DataFrame({'matricule': [matricule],
                           'nom': [nom],
                           'courriel': [courriel]})
        self.append(df)

        cond = self.db.table('personnes').columns.courriel == courriel
        index_participant = self.select(where=[cond])
        return index_participant.index[0]

    def trouver(self, matricule=None, nom=None, courriel=None):
        conds = []

        if matricule is not None:
            cond = self.db.table('personnes').columns.matricule == matricule
            conds.append(cond)
        if nom is not None:
            cond = self.db.table('personnes').columns.nom == nom
            conds.append(cond)
        if courriel is not None:
            cond = self.db.table('personnes').columns.courriel == courriel
            conds.append(cond)

        return self.select(where=conds)


class FormationLaser(BaseTableau):
    colonnes_standard = (column('index', int, primary_key=True),
                         column('personne', int),
                         column('date', dt),
                         column('validation', bool))

    def __init__(self, adresse: str, reflect: bool = True):
        table = 'sst_laser'

        metadata = sqla.MetaData()
        db = BaseDeDonnées(adresse, metadata)

        if reflect:
            moteur = db.create_engine()
            metadata.reflect(moteur)
        else:
            sqla.Table(table,
                       metadata,
                       *self.colonnes_standard)

        super().__init__(db, table)

        self.participants = Participants(adresse, reflect)

    def ajouter(self, matricule, nom, courriel, date=None):
        print('On ajoute une entrée dans la base de données.')
        index_participant = self.participants.ajouter(matricule, nom, courriel)

        if date is None:
            date = dt.today()

        df = pd.DataFrame({'personne': [index_participant],
                           'date': [date],
                           'validation': [False]})
        self.append(df)

    def trouver(self, matricule=None, nom=None, courriel=None, début=None):
        personnes = self.participants.trouver(matricule, nom, courriel).index

        dfs = []
        for personne in personnes:
            conds = []
            cond = self.db.table('sst_laser').columns.index == personne
            conds.append(cond)
            if début is not None:
                cond = self.db.table('sst_laser').columns.date >= début
                conds.append(cond)

            dfs.append(self.select(where=conds))

        return pd.concat(dfs)


class SSTLaserCertificatsForm(MSForm):

    def nettoyer(self, cadre):
        cadre = self.convertir_champs(cadre)
        cadre = cadre.astype({'matricule': int}, errors='ignore')

        courriels_manquants = cadre['courriel'] == 'anonymous'
        cadre.loc[courriels_manquants,
                  'courriel'] = cadre.loc[courriels_manquants, 'courriel2']
        cadre.courriel = cadre.courriel.fillna(
            cadre.courriel2).fillna('@polymtl.ca')
        cadre.nom = cadre.nom.fillna(cadre.nom2).fillna('anonyme')
        cadre.date = cadre.date.dt.date
        cadre.matricule = cadre.matricule.fillna(0)

        return cadre.loc[:, ['date', 'matricule', 'courriel', 'nom', 'nom_responsable', 'courriel_responsable', 'manipulation', 'alignement']]

    def courriel(self, entrée):
        nom_fichier = 'temp.xlsx'
        entrée.to_excel(nom_fichier)

        destinataire = 'emile.jetzer@polymtl.ca'
        sujet = 'Complétion du test de sécurité laser'
        message = f'{entrée.nom} ({entrée.matricule}, {entrée.courriel}), vient de compléter la formation de sécurité laser. Il travaille sous la supervision de {entrée.nom_responsable} ({entrée.courriel_responsable}).'

        if entrée.manipulation == 'Je vais manipuler des lasers de classe 2 et plus'\
                and entrée.alignement == 'Non':
            destinataire += ',mikael.leduc@polymtl.ca'
            message += '\n\nLa formation d\'alignement est nécessaire.'
            sujet = '[formation alignement] ' + sujet

        Courriel(destinataire,
                 'emile.jetzer@polymtl.ca',
                 'Complétion du test de sécurité laser',
                 message,
                 pièces_jointes=(nom_fichier,)).envoyer('smtp.polymtl.ca')

    def action(self, cadre):
        chemin_cert = Path(__file__).parent / self.config.get('certificats', 'chemin')
        cert = Certificat(chemin_cert)
        for i, entrée in cadre.iterrows():
            nom_participant, matricule, courriel, date = entrée.nom, entrée.matricule, entrée.courriel, entrée.date
            cert.màj(nom_participant, matricule)

            self.courriel(entrée)

            for disque in self.config.getlist('certificats', 'disques'):
                url = self.config.get(disque, 'url')
                chemin = self.config.getpath(disque, 'mount_point')
                drive = self.config.get(disque, 'drive')
                mode = self.config.get(disque, 'method')

                nom = self.config.get(disque, 'nom')
                mdp = keyring.get_password(
                    'system', f'polygphys.sst.laser.{disque}.{nom}')
                if mdp is None:
                    mdp = getpass.getpass('mdp: ')
                    keyring.set_password(
                        'system', f'polygphys.sst.laser.{disque}.{nom}', mdp)
                with DisqueRéseau(url, chemin, drive, nom, mdp, mode) as d:
                    sous_dossier = d / self.config.get(disque, 'chemin')
                    sous_dossier = d / self.config.get('certificats', 'ppt')
                    fichier = sous_dossier / f'{entrée.nom}.pptx'
                    cert.enregistrer(fichier)

            adresse = self.config.get('bd', 'adresse')
            base_de_données = FormationLaser(adresse)
            base_de_données.ajouter(matricule, nom_participant, courriel, date)

# Programme


def main():
    import logging
    logging.info('On reste vigilant pour les nouveaux certificats...')
    chemin_config = Path('~').expanduser() / 'certificats_laser.cfg'
    config = SSTLaserCertificatsConfig(chemin_config)

    dossier = OneDrive('',
                       config.get('onedrive', 'organisation'),
                       config.get('onedrive', 'sous-dossier'),
                       partagé=True)
    fichier = dossier / config.get('formulaire', 'nom')
    config.set('formulaire', 'chemin', str(fichier))

    formulaire = SSTLaserCertificatsForm(config)
    logging.info('unoconv démarre.')
    exporteur = subprocess.Popen(['unoconv', '--listener'])

    schedule.every().day.at('08:00').do(formulaire.mise_à_jour)
    
    formulaire.mise_à_jour()
    try:
        while True:
            schedule.run_pending()
            time.sleep(1)
    except KeyboardInterrupt:
        logging.info('On arrête.')
    finally:
        exporteur.terminate()

    logging.info('Terminé.')


if __name__ == '__main__':
    main()
